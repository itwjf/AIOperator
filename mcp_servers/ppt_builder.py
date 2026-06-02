"""
PPT 渲染引擎 — 加载模板、生成页面、保存文件。

所有排版逻辑集中在这里，AI 只传内容数据不碰样式。
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
import os

# 模板路径
TEMPLATE_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "templates")
DEFAULT_TEMPLATE = os.path.join(TEMPLATE_DIR, "default.pptx")

# 页面尺寸常量
SLIDE_WIDTH = Inches(13.333)  # 16:9 宽屏
LEFT_MARGIN = Inches(1.0)
CONTENT_WIDTH = Inches(11.333)


class PPTBuilder:
    """PPT 构建器 — 持有当前编辑的 Presentation 对象。"""

    def __init__(self, template_path: str = DEFAULT_TEMPLATE):
        if os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
            self.prs.slide_width = SLIDE_WIDTH
        self.slide_count = 0

    # === 页面生成 ===

    def add_cover(self, title: str, subtitle: str = "", author: str = "") -> int:
        """生成封面页。返回页码（1-based）。"""
        layout = self.prs.slide_layouts[0]  # 第一个布局通常适合做封面
        slide = self.prs.slides.add_slide(layout)

        # 标题
        if slide.shapes.title:
            slide.shapes.title.text = title

        # 副标题（如果有 placeholder）
        if subtitle:
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1 and shape is not slide.shapes.title:
                    shape.text = subtitle
                    break

        self.slide_count += 1
        return self.slide_count

    def add_content_slide(self, title: str, bullets: list[str]) -> int:
        """生成文字要点页。返回页码。"""
        layout = self.prs.slide_layouts[1]  # 标题+内容布局
        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        # 要点文本填充到内容 placeholder
        body = ""
        for b in bullets:
            body += f"• {b}\n"

        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1 and shape is not slide.shapes.title:
                shape.text = body
                break

        self.slide_count += 1
        return self.slide_count

    def add_table_slide(self, title: str, columns: list[str], rows: list[list]) -> int:
        """生成表格页。列数和行数完全动态，自动适配列宽和字号。"""
        layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(layout)

        # 标题
        if slide.shapes.title:
            slide.shapes.title.text = title

        n_rows = len(rows) + 1  # +1 for header
        n_cols = len(columns)

        # 自适应参数
        col_width = CONTENT_WIDTH / n_cols
        font_size = Pt(8) if n_cols > 6 else Pt(11)

        # 表格位置
        left = LEFT_MARGIN
        top = Inches(1.8)
        table_height = Inches(min(n_rows * 0.35, 5.5))

        table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, CONTENT_WIDTH, table_height)
        table = table_shape.table

        # 写表头
        for c, col_name in enumerate(columns):
            cell = table.cell(0, c)
            cell.text = str(col_name)
            for p in cell.text_frame.paragraphs:
                p.font.size = font_size
                p.font.bold = True

        # 写数据行
        for r, row_data in enumerate(rows):
            for c, value in enumerate(row_data):
                cell = table.cell(r + 1, c)
                cell.text = str(value) if value is not None else ""
                for p in cell.text_frame.paragraphs:
                    p.font.size = font_size

        self.slide_count += 1
        return self.slide_count

    def save(self, output_path: str) -> str:
        """保存 pptx 文件，返回绝对路径。"""
        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        self.prs.save(output_path)
        return os.path.abspath(output_path)


# === 会话管理 ===
# 同一会话多次工具调用共享同一个 PPTBuilder 实例
_sessions: dict[str, PPTBuilder] = {}


def get_builder(session_id: str) -> PPTBuilder:
    """获取或创建指定会话的 PPTBuilder。"""
    if session_id not in _sessions:
        _sessions[session_id] = PPTBuilder()
    return _sessions[session_id]
